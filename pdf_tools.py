"""
PDF Tools — the processing engine that runs entirely in the browser via Pyodide.

Loaded into a Web Worker; every function operates on in-memory bytes handed over
from the page. Nothing here opens a socket or touches a server — the user's file
never leaves their machine.

Public entry point is `dispatch(action, params_json)`; inputs are read from the
Pyodide virtual FS at /in0, /in1, ... and binary results are written to /out.
"""

import io
import csv
import json
import logging
import re

from pypdf import PdfReader, PdfWriter

# ---- progress streaming ----------------------------------------------------
# The worker installs a JS callback via set_progress(); long-running engines
# (pdf2docx logs one line per page) then stream real progress to the UI.
_PROGRESS = None


def set_progress(cb) -> None:
    global _PROGRESS
    _PROGRESS = cb


_ANSI = re.compile(r"\x1b\[[0-9;]*[A-Za-z]")  # pdf2docx colors its log lines


class _ProgressHandler(logging.Handler):
    def emit(self, record):
        if _PROGRESS is not None:
            try:
                _PROGRESS(_ANSI.sub("", record.getMessage())[:90])
            except Exception:
                pass


logging.getLogger().addHandler(_ProgressHandler())
logging.getLogger().setLevel(logging.INFO)


# --------------------------------------------------------------------------- #
# helpers
# --------------------------------------------------------------------------- #
def _reader(data: bytes) -> PdfReader:
    return PdfReader(io.BytesIO(data))


def _save(writer: PdfWriter) -> bytes:
    out = io.BytesIO()
    writer.write(out)
    return out.getvalue()


def _latin1(s) -> str:
    # fpdf2 core fonts are latin-1 only; used only when the font pack is absent.
    return str(s).encode("latin-1", "replace").decode("latin-1")


FONT_DIR = "/fonts"


def _register_fonts(pdf):
    """Register the bundled Unicode fonts (worker writes them into /fonts) on an
    FPDF instance. Returns the body family — NotoSans with a Devanagari fallback
    when the pack is present, or None to signal core-font (latin-1) mode."""
    import os

    if not os.path.isfile(f"{FONT_DIR}/NotoSans-Regular.ttf"):
        return None
    try:
        pdf.add_font("NotoSans", "", f"{FONT_DIR}/NotoSans-Regular.ttf")
        pdf.add_font("NotoSans", "B", f"{FONT_DIR}/NotoSans-Bold.ttf")
        pdf.add_font("NotoDeva", "", f"{FONT_DIR}/NotoSansDevanagari-Regular.ttf")
        if os.path.isfile(f"{FONT_DIR}/NotoSansDevanagari-Bold.ttf"):
            pdf.add_font("NotoDeva", "B", f"{FONT_DIR}/NotoSansDevanagari-Bold.ttf")
        try:
            pdf.set_fallback_fonts(["NotoDeva"], exact_match=False)
        except TypeError:
            pdf.set_fallback_fonts(["NotoDeva"])
        return "NotoSans"
    except Exception:
        return None


def _face(pdf) -> str:
    return getattr(pdf, "uni_family", None) or "Helvetica"


def _txt(pdf, s) -> str:
    return str(s) if getattr(pdf, "uni_family", None) else _latin1(s)


def parse_page_ranges(spec: str, total: int):
    """'2, 4, 7-9' -> sorted list of 1-indexed pages within [1, total]."""
    pages = set()
    for part in (spec or "").split(","):
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            a, _, b = part.partition("-")
            try:
                lo, hi = int(a), int(b)
            except ValueError:
                continue
            lo, hi = min(lo, hi), max(lo, hi)
            for p in range(lo, hi + 1):
                if 1 <= p <= total:
                    pages.add(p)
        else:
            try:
                p = int(part)
            except ValueError:
                continue
            if 1 <= p <= total:
                pages.add(p)
    return sorted(pages)


def get_page_count(data: bytes) -> int:
    return len(_reader(data).pages)


# --------------------------------------------------------------------------- #
# core PDF operations (pypdf)
# --------------------------------------------------------------------------- #
def merge_pdfs(files) -> bytes:
    if not files:
        raise ValueError("No files to merge.")
    writer = PdfWriter()
    for data in files:
        writer.append(_reader(data))
    return _save(writer)


def split_pdf(data: bytes, start: int, end: int) -> bytes:
    reader = _reader(data)
    total = len(reader.pages)
    if not (1 <= start <= total):
        raise ValueError(f"Start page must be between 1 and {total}.")
    if not (start <= end <= total):
        raise ValueError(f"End page must be between {start} and {total}.")
    writer = PdfWriter()
    for i in range(start - 1, end):
        writer.add_page(reader.pages[i])
    return _save(writer)


def delete_pages(data: bytes, spec: str) -> bytes:
    reader = _reader(data)
    total = len(reader.pages)
    remove = set(parse_page_ranges(spec, total))
    if not remove:
        raise ValueError("Enter at least one valid page to remove.")
    if len(remove) >= total:
        raise ValueError("At least one page must remain.")
    writer = PdfWriter()
    for i in range(total):
        if (i + 1) not in remove:
            writer.add_page(reader.pages[i])
    return _save(writer)


def organize_pdf(data: bytes, order) -> bytes:
    """order: list of {"src": 0-based index, "rot": added clockwise degrees}."""
    reader = _reader(data)
    total = len(reader.pages)
    writer = PdfWriter()
    for item in order:
        src = int(item["src"])
        if not (0 <= src < total):
            continue
        page = reader.pages[src]
        rot = int(item.get("rot", 0)) % 360
        if rot:
            page.rotate(rot)
        writer.add_page(page)
    if len(writer.pages) == 0:
        raise ValueError("The document needs at least one page.")
    return _save(writer)


def compress_pdf(data: bytes, image_quality: int = 60) -> bytes:
    """Lossless stream compression + duplicate-object removal, plus optional
    image recompression. Always returns a VALID pdf (never truncated)."""
    reader = _reader(data)
    writer = PdfWriter()
    writer.append(reader)

    for page in writer.pages:
        try:
            page.compress_content_streams()
        except Exception:
            pass

    if image_quality and image_quality < 100:
        for page in writer.pages:
            try:
                for img in page.images:
                    img.replace(img.image, quality=image_quality)
            except Exception:
                pass

    try:
        writer.compress_identical_objects()
    except Exception:
        pass

    return _save(writer)


# --------------------------------------------------------------------------- #
# overlays: watermark / page numbers / invisible OCR text (fpdf2 + pypdf)
# --------------------------------------------------------------------------- #
def _overlay_merge(data: bytes, draw) -> bytes:
    """Draw per-page content with fpdf2 and merge it on top of the original
    pages. draw(ov, i, w, h) gets the overlay canvas and the page size in pt."""
    from fpdf import FPDF

    reader = _reader(data)
    ov = FPDF(unit="pt")
    ov.set_auto_page_break(False)
    ov.uni_family = _register_fonts(ov)
    for page in reader.pages:
        w, h = float(page.mediabox.width), float(page.mediabox.height)
        ov.add_page(format=(w, h))
        draw(ov, ov.page - 1, w, h)
    ov_reader = _reader(bytes(ov.output()))
    writer = PdfWriter()
    for i, page in enumerate(reader.pages):
        page.merge_page(ov_reader.pages[i])
        writer.add_page(page)
    return _save(writer)


def stamp_pdf(data: bytes, text: str = "", pos: str = "diagonal",
              pagenum: bool = False) -> bytes:
    """Watermark / header / footer text, and/or 'Page i of N' numbers."""
    if not text.strip() and not pagenum:
        raise ValueError("Add stamp text or turn on page numbers.")
    total = get_page_count(data)

    def draw(ov, i, w, h):
        face = _face(ov)
        if text.strip():
            t = _txt(ov, text.strip())
            if pos == "diagonal":
                ov.set_font(face, "B", max(24, min(w, h) / 8))
                ov.set_text_color(203, 178, 130)
                with ov.rotation(45, w / 2, h / 2):
                    ov.text((w - ov.get_string_width(t)) / 2, h / 2, t)
            else:
                ov.set_font(face, "", 9)
                ov.set_text_color(122, 110, 92)
                y = 24 if pos == "header" else h - 16
                ov.text((w - ov.get_string_width(t)) / 2, y, t)
        if pagenum:
            ov.set_font(face, "", 9)
            ov.set_text_color(122, 110, 92)
            label = f"Page {i + 1} of {total}"
            ov.text(w - ov.get_string_width(label) - 28, h - 16, label)

    return _overlay_merge(data, draw)


def ocr_overlay(data: bytes, pages) -> bytes:
    """Lay recognized words invisibly over each page → a searchable PDF that
    still looks exactly like the scan. pages: [{scale, words: [{t, x0, y0,
    x1, y1}]}] with word boxes in canvas pixels at the given render scale."""
    def draw(ov, i, w, h):
        if i >= len(pages):
            return
        s = float(pages[i].get("scale") or 1) or 1
        words = pages[i].get("words") or []
        if not words:
            return
        face = _face(ov)
        try:
            from fpdf.enums import TextMode
            ov.text_mode = TextMode.INVISIBLE
        except Exception:
            ov.set_text_color(255, 255, 255)
        for wd in words:
            t = str(wd.get("t", "")).strip()
            if not t:
                continue
            x0, y0, y1 = wd["x0"] / s, wd["y0"] / s, wd["y1"] / s
            size = max(4.0, (y1 - y0) * 0.92)
            ov.set_font(face, "", size)
            ov.text(x0, y1 - size * 0.22, _txt(ov, t))

    return _overlay_merge(data, draw)


def strip_metadata(data: bytes) -> bytes:
    """A privacy pass: drop the document info dictionary and XMP metadata."""
    reader = _reader(data)
    writer = PdfWriter()
    writer.append(reader)
    try:
        writer.metadata = None
    except Exception:
        writer.add_metadata({})
    try:
        root = writer._root_object
        if "/Metadata" in root:
            del root["/Metadata"]
    except Exception:
        pass
    return _save(writer)


def split_zip(data: bytes, mode: str = "every", n: int = 0, spec: str = "") -> bytes:
    """Split into several PDFs delivered as one zip: every N pages, or a list
    of ranges ('1-3, 7, 9-12' → one PDF per comma-separated chunk)."""
    import zipfile

    reader = _reader(data)
    total = len(reader.pages)
    parts = []
    if mode == "every":
        n = max(1, int(n))
        for s in range(0, total, n):
            e = min(s + n, total)
            parts.append((f"pages_{s + 1:03d}-{e:03d}", range(s, e)))
    else:
        for chunk in (spec or "").split(","):
            chunk = chunk.strip()
            if not chunk:
                continue
            pages = parse_page_ranges(chunk, total)
            if pages:
                parts.append((f"pages_{chunk.replace(' ', '').replace('-', '_')}",
                              [p - 1 for p in pages]))
        if not parts:
            raise ValueError("Enter ranges like 1-3, 7, 9-12.")

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as z:
        for name, idxs in parts:
            w = PdfWriter()
            for i in idxs:
                w.add_page(reader.pages[i])
            b = io.BytesIO()
            w.write(b)
            z.writestr(name + ".pdf", b.getvalue())
    return buf.getvalue()


def compress_pdf_max(data: bytes, dpi: int = 110, quality: int = 60) -> bytes:
    """Ghostscript-style compression via PyMuPDF (worker loads it on demand):
    downsample images above ~1.3× the target DPI, recompress as JPEG, then
    rebuild with full garbage collection. Never returns a larger file."""
    import pymupdf

    doc = pymupdf.open(stream=data, filetype="pdf")
    try:
        doc.rewrite_images(dpi_threshold=int(dpi * 1.3), dpi_target=dpi,
                           quality=quality, lossy=True, lossless=True)
        out = doc.tobytes(garbage=4, deflate=True)
    finally:
        doc.close()
    return out if len(out) < len(data) else data


# --------------------------------------------------------------------------- #
# text extraction / generation
# --------------------------------------------------------------------------- #
def pdf_to_text(data: bytes) -> str:
    reader = _reader(data)
    chunks = []
    for i, page in enumerate(reader.pages, 1):
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""
        chunks.append(f"--- Page {i} ---\n{text.strip()}")
    return "\n\n".join(chunks).strip() or "[No selectable text found in this PDF.]"


def _new_pdf(orientation="P", fmt="letter"):
    from fpdf import FPDF

    pdf = FPDF(orientation=orientation, unit="pt", format=fmt)
    pdf.set_auto_page_break(True, margin=48)
    pdf.uni_family = _register_fonts(pdf)
    return pdf


def _mcell(pdf, h, text):
    """multi_cell that survives unbreakable long tokens (URLs, base64, …)."""
    text = text if text.strip() else " "
    try:
        pdf.multi_cell(0, h, text, new_x="LMARGIN", new_y="NEXT")
    except Exception:
        pdf.multi_cell(0, h, text, new_x="LMARGIN", new_y="NEXT", wrapmode="CHAR")


def text_to_pdf(text: str, name: str = "document") -> bytes:
    pdf = _new_pdf()
    pdf.add_page()
    pdf.set_font(_face(pdf), "B", 9)
    pdf.set_text_color(120, 120, 120)
    pdf.cell(0, 12, _txt(pdf, name), new_x="LMARGIN", new_y="NEXT")
    pdf.ln(6)
    pdf.set_font(_face(pdf), size=11)
    pdf.set_text_color(20, 20, 20)
    for line in _txt(pdf, text).split("\n"):
        _mcell(pdf, 15, line)
    return bytes(pdf.output())


# --------------------------------------------------------------------------- #
# spreadsheets (openpyxl) <-> PDF (fpdf2)
# --------------------------------------------------------------------------- #
def _rows_from_csv(data: bytes):
    text = data.decode("utf-8", errors="replace")
    return list(csv.reader(io.StringIO(text)))


def _cellstr(v) -> str:
    import datetime as dt

    if v is None:
        return ""
    if isinstance(v, dt.datetime):
        return v.strftime("%Y-%m-%d %H:%M") if (v.hour or v.minute or v.second) else v.strftime("%Y-%m-%d")
    if isinstance(v, dt.date):
        return v.strftime("%Y-%m-%d")
    if isinstance(v, float) and v == int(v):
        return str(int(v))
    return str(v)


def _argb(color):
    """openpyxl color → (r, g, b), or None for absent/transparent/theme colors."""
    try:
        rgb = getattr(color, "rgb", None)
        if isinstance(rgb, str) and len(rgb) == 8 and rgb[:2] != "00":
            return tuple(int(rgb[i:i + 2], 16) for i in (2, 4, 6))
    except Exception:
        pass
    return None


def _xlsx_to_pdf(data: bytes) -> bytes:
    """Styled, multi-sheet XLSX rendering: one section per sheet, real column
    widths, bold/fill/font colors, alignment, dates formatted, horizontal
    merges honoured via colspan. (Italic is skipped — the bundled Unicode
    family ships no italic face.)"""
    from fpdf import FPDF
    from fpdf.fonts import FontFace
    from openpyxl import load_workbook
    from openpyxl.utils import get_column_letter

    wb = load_workbook(io.BytesIO(data), data_only=True)
    pdf = FPDF(orientation="L", unit="pt", format="a4")
    pdf.set_auto_page_break(True, margin=32)
    pdf.uni_family = _register_fonts(pdf)
    face = _face(pdf)

    rendered = 0
    for ws in wb.worksheets:
        nrows = min(ws.max_row or 0, 3000)
        ncols = min(ws.max_column or 0, 60)
        if not nrows or not ncols or (nrows == 1 and ncols == 1 and ws.cell(1, 1).value is None):
            continue
        rendered += 1

        pdf.add_page()
        pdf.set_font(face, "B", 12)
        pdf.set_text_color(74, 56, 34)
        pdf.cell(0, 16, _txt(pdf, ws.title), new_x="LMARGIN", new_y="NEXT")
        pdf.ln(6)
        pdf.set_font(face, size=8)
        pdf.set_text_color(20, 20, 20)

        h_merge = {}  # (row, col) → colspan for single-row merges
        for rng in ws.merged_cells.ranges:
            if rng.min_row == rng.max_row and rng.max_col > rng.min_col:
                h_merge[(rng.min_row, rng.min_col)] = min(rng.max_col, ncols) - rng.min_col + 1

        widths = []
        for ci in range(1, ncols + 1):
            dim = ws.column_dimensions.get(get_column_letter(ci))
            w = dim.width if dim is not None and dim.width else 10.5
            widths.append(max(4.0, min(40.0, float(w))))  # relative proportions

        with pdf.table(col_widths=tuple(widths), first_row_as_headings=False,
                       line_height=11) as table:
            for ri in range(1, nrows + 1):
                trow = table.row()
                ci = 1
                while ci <= ncols:
                    cell = ws.cell(ri, ci)
                    span = h_merge.get((ri, ci), 1)
                    style = FontFace(
                        emphasis="BOLD" if (cell.font is not None and cell.font.b) else None,
                        color=_argb(cell.font.color) if cell.font is not None else None,
                        fill_color=_argb(cell.fill.fgColor)
                        if cell.fill is not None and cell.fill.patternType == "solid" else None,
                    )
                    halign = (cell.alignment.horizontal or "LEFT").upper() \
                        if cell.alignment is not None and cell.alignment.horizontal in ("center", "right") else "LEFT"
                    trow.cell(_txt(pdf, _cellstr(cell.value)), colspan=span,
                              style=style, align=halign)
                    ci += span

    if not rendered:
        raise ValueError("The workbook has no visible content.")
    return bytes(pdf.output())


def table_to_pdf(data: bytes, is_csv: bool) -> bytes:
    if not is_csv:
        return _xlsx_to_pdf(data)
    from fpdf import FPDF

    rows = _rows_from_csv(data)
    rows = [r for r in rows if any(str(c).strip() for c in r)] or [["(empty file)"]]
    ncols = max(len(r) for r in rows)

    pdf = FPDF(orientation="L", unit="pt", format="a4")
    pdf.set_auto_page_break(True, margin=32)
    pdf.uni_family = _register_fonts(pdf)
    rows = [[_txt(pdf, c) for c in r] + [""] * (ncols - len(r)) for r in rows]
    pdf.add_page()
    pdf.set_font(_face(pdf), size=8)
    with pdf.table(first_row_as_headings=True) as table:
        for r in rows:
            trow = table.row()
            for cell in r:
                trow.cell(cell)
    return bytes(pdf.output())


def pdf_to_xlsx(data: bytes) -> bytes:
    from openpyxl import Workbook

    reader = _reader(data)
    wb = Workbook()
    ws = wb.active
    ws.title = "Extracted"
    for i, page in enumerate(reader.pages, 1):
        try:
            text = page.extract_text(extraction_mode="layout") or ""
        except Exception:
            text = page.extract_text() or ""
        for line in text.split("\n"):
            if not line.strip():
                continue
            import re
            cells = re.split(r"\s{2,}", line.strip())
            ws.append(cells)
        if i < len(reader.pages):
            ws.append([f"— page {i} / {i + 1} —"])
    out = io.BytesIO()
    wb.save(out)
    return out.getvalue()


# --------------------------------------------------------------------------- #
# Word (python-docx) <-> PDF
# --------------------------------------------------------------------------- #
def word_to_pdf(data: bytes) -> bytes:
    import docx

    doc = docx.Document(io.BytesIO(data))
    pdf = _new_pdf()
    pdf.add_page()
    face = _face(pdf)
    for para in doc.paragraphs:
        text = _txt(pdf, para.text)
        style = (para.style.name or "").lower() if para.style else ""
        if not text.strip():
            pdf.ln(8)
            continue
        if "heading 1" in style or style == "title":
            pdf.set_font(face, "B", 18)
        elif "heading" in style:
            pdf.set_font(face, "B", 14)
        else:
            pdf.set_font(face, size=11)
        _mcell(pdf, 16, text)
        pdf.ln(2)
    return bytes(pdf.output())


def pdf_to_word_hifi(data: bytes, start: int = 0, end: int = 0) -> bytes:
    """Layout-aware PDF→DOCX via pdf2docx: flowing paragraphs, ruled tables,
    images. The worker installs the engine on first use. start/end are
    1-based inclusive page numbers (0 = from first / to last) — converting a
    range is the practical answer for very large documents."""
    from pdf2docx import Converter

    with open("/hifi.pdf", "wb") as fh:
        fh.write(data)
    kwargs = {}
    if start:
        kwargs["start"] = int(start) - 1   # pdf2docx: 0-based inclusive
    if end:
        kwargs["end"] = int(end)           # pdf2docx: 0-based exclusive
    cv = Converter("/hifi.pdf")
    try:
        cv.convert("/hifi.docx", **kwargs)
    finally:
        cv.close()
    with open("/hifi.docx", "rb") as fh:
        return fh.read()


def pdf_to_xlsx_hifi(data: bytes) -> bytes:
    """Layout-aware PDF→XLSX via pdfplumber's table detection; pages without
    detectable tables fall back to whitespace-split text rows."""
    import re

    import pdfplumber
    from openpyxl import Workbook

    wb = Workbook()
    wb.remove(wb.active)
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            ws = wb.create_sheet(f"Page {i}"[:31])
            tables = page.extract_tables()
            if tables:
                for t in tables:
                    for row in t:
                        ws.append(["" if c is None else str(c) for c in row])
                    ws.append([])
            else:
                for line in (page.extract_text() or "").split("\n"):
                    if line.strip():
                        ws.append(re.split(r"\s{2,}", line.strip()))
    if not wb.sheetnames:
        wb.create_sheet("Empty")
    out = io.BytesIO()
    wb.save(out)
    return out.getvalue()


def pdf_to_word(data: bytes) -> bytes:
    import docx

    reader = _reader(data)
    doc = docx.Document()
    for i, page in enumerate(reader.pages, 1):
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""
        for block in text.split("\n"):
            doc.add_paragraph(block)
        if i < len(reader.pages):
            doc.add_page_break()
    out = io.BytesIO()
    doc.save(out)
    return out.getvalue()


# --------------------------------------------------------------------------- #
# PowerPoint (python-pptx) <-> PDF
# --------------------------------------------------------------------------- #
def _shape_texts(shape, parts):
    """Collect text from a shape, recursing into groups and reading tables —
    designed decks keep most of their words in exactly those two places."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub in shape.shapes:
            _shape_texts(sub, parts)
        return
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            if para.text.strip():
                parts.append(para.text)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append("  ·  ".join(c for c in cells if c))


def _slide_texts(prs):
    slides = []
    for slide in prs.slides:
        parts = []
        for shape in slide.shapes:
            _shape_texts(shape, parts)
        slides.append(parts)
    return slides


def ppt_to_pdf(data: bytes) -> bytes:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    texts = _slide_texts(prs)
    if not any(texts):
        raise ValueError("No extractable text in this presentation — its slides "
                         "appear to be images or artwork. The current converter "
                         "re-lays text only; a pixel-faithful one is on the roadmap.")
    pdf = _new_pdf(orientation="L", fmt="a4")
    face = _face(pdf)
    for parts in texts:
        pdf.add_page()
        title = parts[0] if parts else "Slide"
        body = parts[1:] if len(parts) > 1 else []
        pdf.set_font(face, "B", 22)
        _mcell(pdf, 28, _txt(pdf, title))
        pdf.ln(8)
        pdf.set_font(face, size=13)
        for line in body:
            _mcell(pdf, 20, _txt(pdf, "-  " + line))
    if pdf.page_no() == 0:
        pdf.add_page()
    return bytes(pdf.output())


def pdf_to_ppt(data: bytes) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    reader = _reader(data)
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for i, page in enumerate(reader.pages, 1):
        try:
            text = (page.extract_text() or "").strip()
        except Exception:
            text = ""
        slide = prs.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Inches(0.6), Inches(0.5),
                                       prs.slide_width - Inches(1.2),
                                       prs.slide_height - Inches(1.0))
        tf = box.text_frame
        tf.word_wrap = True
        tf.text = f"Page {i}"
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.bold = True
        for line in text.split("\n"):
            if line.strip():
                p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(12)
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def images_to_ppt(paths, w_pt: float = 0, h_pt: float = 0) -> bytes:
    """Faithful PDF→PPTX: each pre-rendered page image becomes a full-bleed
    slide picture. Visually exact; the text is not editable — that trade is
    stated in the UI."""
    from pptx import Presentation
    from pptx.util import Emu

    EMU_PER_PT = 12700
    prs = Presentation()
    if w_pt and h_pt:
        prs.slide_width = Emu(int(w_pt * EMU_PER_PT))
        prs.slide_height = Emu(int(h_pt * EMU_PER_PT))
    blank = prs.slide_layouts[6]
    for path in paths:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(path, 0, 0,
                                 width=prs.slide_width, height=prs.slide_height)
    if not prs.slides:
        prs.slides.add_slide(blank)
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


# --------------------------------------------------------------------------- #
# images -> PDF (Pillow + fpdf2)
# --------------------------------------------------------------------------- #
def images_to_pdf(paths, size: str = "letter") -> bytes:
    from fpdf import FPDF
    from PIL import Image

    dims = {"letter": (612, 792), "a4": (595, 842)}
    pw, ph = dims.get(size, dims["letter"])
    margin = 24
    pdf = FPDF(unit="pt", format=(pw, ph))
    for path in paths:
        with Image.open(path) as im:
            iw, ih = im.size
        avail_w, avail_h = pw - 2 * margin, ph - 2 * margin
        scale = min(avail_w / iw, avail_h / ih)
        w, h = iw * scale, ih * scale
        pdf.add_page()
        pdf.image(path, x=(pw - w) / 2, y=(ph - h) / 2, w=w, h=h)
    return bytes(pdf.output())


# --------------------------------------------------------------------------- #
# dispatch
# --------------------------------------------------------------------------- #
def _inputs(n):
    files = []
    for i in range(n):
        with open(f"/in{i}", "rb") as f:
            files.append(f.read())
    return files


def dispatch(action: str, params_json: str = "") -> str:
    p = json.loads(params_json) if params_json else {}
    n = int(p.get("n", 1))

    # metadata / text results return JSON directly (no /out file)
    if action == "pageCount":
        return json.dumps({"kind": "json", "pages": get_page_count(_inputs(1)[0])})
    if action == "pdfToText":
        return json.dumps({"kind": "text", "text": pdf_to_text(_inputs(1)[0])})

    files = _inputs(n)
    if action == "merge":
        out = merge_pdfs(files)
    elif action == "split":
        out = split_pdf(files[0], int(p["start"]), int(p["end"]))
    elif action == "delete":
        out = delete_pages(files[0], p["ranges"])
    elif action == "organize":
        out = organize_pdf(files[0], p["order"])
    elif action == "stamp":
        out = stamp_pdf(files[0], p.get("text", ""), p.get("pos", "diagonal"),
                        bool(p.get("pagenum")))
    elif action == "stripMeta":
        out = strip_metadata(files[0])
    elif action == "splitZip":
        out = split_zip(files[0], p.get("mode", "every"), int(p.get("n", 0)),
                        p.get("spec", ""))
    elif action == "ocrOverlay":
        out = ocr_overlay(files[0], p.get("pages") or [])
    elif action == "compress":
        if p.get("mode") == "max":
            out = compress_pdf_max(files[0], int(p.get("dpi", 110)),
                                   int(p.get("quality", 60)))
        else:
            out = compress_pdf(files[0], int(p.get("quality", 60)))
    elif action == "textToPdf":
        out = text_to_pdf(p.get("text", ""), p.get("name", "document"))
    elif action == "tableToPdf":
        out = table_to_pdf(files[0], bool(p.get("isCsv")))
    elif action == "pdfToXlsx":
        out = pdf_to_xlsx_hifi(files[0]) if p.get("engine") == "hifi" else pdf_to_xlsx(files[0])
    elif action == "wordToPdf":
        out = word_to_pdf(files[0])
    elif action == "pdfToWord":
        out = (pdf_to_word_hifi(files[0], int(p.get("start", 0)), int(p.get("end", 0)))
               if p.get("engine") == "hifi" else pdf_to_word(files[0]))
    elif action == "pptToPdf":
        out = ppt_to_pdf(files[0])
    elif action == "pdfToPpt":
        out = pdf_to_ppt(files[0])
    elif action == "imagesToPdf":
        for i, data in enumerate(files):
            with open(f"/img{i}", "wb") as f:
                f.write(data)
        out = images_to_pdf([f"/img{i}" for i in range(len(files))],
                            p.get("size", "letter"))
    elif action == "imagesToPpt":
        for i, data in enumerate(files):
            with open(f"/img{i}", "wb") as f:
                f.write(data)
        out = images_to_ppt([f"/img{i}" for i in range(len(files))],
                            float(p.get("w", 0)), float(p.get("h", 0)))
    else:
        raise ValueError(f"Unknown action: {action}")

    with open("/out", "wb") as f:
        f.write(out)
    return json.dumps({"kind": "file"})
